#!/usr/bin/env python3
"""Generate MCP Governance project presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # Deep navy
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)       # Card background
PURPLE = RGBColor(0x63, 0x66, 0xF1)        # Primary indigo
PURPLE_LIGHT = RGBColor(0xA8, 0x55, 0xF7)  # AI purple
PINK = RGBColor(0xEC, 0x48, 0x99)          # Pink accent
CYAN = RGBColor(0x06, 0xB6, 0xD4)          # Cyan accent
GREEN = RGBColor(0x10, 0xB9, 0x81)         # Green
AMBER = RGBColor(0xF5, 0x9E, 0x0B)         # Amber/warning
RED = RGBColor(0xEF, 0x44, 0x44)           # Red/critical
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=PURPLE_LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"
        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"
    return txBox


def add_icon_card(slide, left, top, width, height, icon, title, desc,
                  accent_color=PURPLE):
    # Card background
    card = add_shape(slide, left, top, width, height, BG_CARD, corner_radius=0.05)
    # Accent bar at top
    add_shape(slide, left, top, width, Inches(0.04), accent_color)
    # Icon
    add_text_box(slide, left + Inches(0.3), top + Inches(0.25), Inches(1), Inches(0.5),
                 icon, font_size=28, color=accent_color, bold=True)
    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.3), top + Inches(1.05), width - Inches(0.6), height - Inches(1.2),
                 desc, font_size=11, color=GRAY)


def add_gradient_accent(slide, left, top, width, height):
    """Add a thin gradient-like accent bar."""
    bar_width = width // 3
    add_shape(slide, left, top, bar_width, height, PURPLE)
    add_shape(slide, left + bar_width, top, bar_width, height, PURPLE_LIGHT)
    add_shape(slide, left + bar_width * 2, top, bar_width + Emu(10), height, PINK)


def add_section_header(slide, text, subtitle=None):
    """Add a consistent section header with accent bar."""
    add_gradient_accent(slide, Inches(0.8), Inches(0.7), Inches(2), Inches(0.05))
    add_text_box(slide, Inches(0.8), Inches(0.85), Inches(8), Inches(0.7),
                 text, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.45), Inches(8), Inches(0.5),
                     subtitle, font_size=16, color=GRAY)


# ═══════════════════════════════════════════════════════
# CREATE PRESENTATION
# ═══════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ══════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# Decorative circles
for cx, cy, r, c, op in [
    (Inches(11), Inches(1), Inches(3), PURPLE, 0.08),
    (Inches(1), Inches(5), Inches(2.5), PURPLE_LIGHT, 0.06),
    (Inches(9), Inches(5.5), Inches(1.5), PINK, 0.05),
]:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    circle.fill.solid()
    circle.fill.fore_color.rgb = c
    circle.line.fill.background()

# Shield icon (text-based)
add_text_box(slide, Inches(5.5), Inches(0.8), Inches(2.5), Inches(1.2),
             "🛡️", font_size=64, alignment=PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(1.5), Inches(2.1), Inches(10.5), Inches(1),
             "MCP Governance", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# AI-Powered badge
badge = add_shape(slide, Inches(5.15), Inches(3.2), Inches(3), Inches(0.5), PURPLE_LIGHT, corner_radius=0.15)
add_text_box(slide, Inches(5.15), Inches(3.22), Inches(3), Inches(0.5),
             "🧠  AI-Powered", font_size=22, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2), Inches(4.0), Inches(9.5), Inches(0.8),
             "Kubernetes-native governance for MCP (Model Context Protocol) infrastructure",
             font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Tech badges row
badges = ["☸ Kubernetes", "🐹 Go 1.25", "⚛ Next.js 14", "🤖 Google Gemini", "🦙 Ollama", "📦 Helm 3"]
badge_width = Inches(1.7)
start_x = Inches(1.5)
for i, badge_text in enumerate(badges):
    bx = start_x + i * (badge_width + Inches(0.15))
    add_shape(slide, bx, Inches(5.2), badge_width, Inches(0.45), BG_CARD, corner_radius=0.1)
    add_text_box(slide, bx, Inches(5.22), badge_width, Inches(0.45),
                 badge_text, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Footer
add_text_box(slide, Inches(3), Inches(6.5), Inches(7.5), Inches(0.5),
             "github.com/techwithhuz/mcp-security-governance",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM — WHY
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Why MCP Governance?", "The Problem: AI Agents Are Expanding — Without Guardrails")

# Problem cards
problems = [
    ("🔓", "No Authentication", "MCP servers exposed without\nJWT or mTLS authentication", RED),
    ("🌐", "No Gateway Control", "Direct access to MCP tools\nwithout proxy enforcement", RED),
    ("🔑", "No Authorization", "No RBAC policies controlling\nwhich agents access which tools", AMBER),
    ("📡", "No Encryption", "MCP traffic flowing without\nTLS encryption", AMBER),
    ("🧰", "Tool Sprawl", "Servers exposing too many tools\nincreasing attack surface", AMBER),
    ("👁️", "No Visibility", "No dashboard or scoring to\nunderstand security posture", PURPLE_LIGHT),
]

for i, (icon, title, desc, color) in enumerate(problems):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(2.2) + row * Inches(2.5)
    add_icon_card(slide, x, y, Inches(3.7), Inches(2.1), icon, title, desc, color)

# Bottom stat
add_shape(slide, Inches(3), Inches(6.7), Inches(7.5), Inches(0.55), BG_CARD, corner_radius=0.08)
add_text_box(slide, Inches(3), Inches(6.72), Inches(7.5), Inches(0.55),
             "⚠️  Without governance, every MCP server is a potential security gap",
             font_size=15, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 3: THE VISION — WHY (continued)
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "The Vision", "A unified governance layer for the MCP ecosystem")

# Vision pillars
pillars = [
    ("🔍", "Discover", "Automatically find all MCP\nresources across every\nnamespace in your cluster", CYAN),
    ("📏", "Evaluate", "Score security posture with\nconfigurable policies across\n8 governance categories", PURPLE),
    ("🧠", "Analyze", "AI-powered deep risk analysis\nwith reasoning, risks, and\nactionable suggestions", PURPLE_LIGHT),
    ("📊", "Visualize", "Real-time enterprise dashboard\nwith trends, findings, and\nper-resource drill-down", GREEN),
]

for i, (icon, title, desc, color) in enumerate(pillars):
    x = Inches(0.6) + i * Inches(3.15)
    # Card
    add_icon_card(slide, x, Inches(2.3), Inches(2.85), Inches(2.8), icon, title, desc, color)
    # Number badge
    num_badge = add_shape(slide, x + Inches(2.1), Inches(2.2), Inches(0.55), Inches(0.55), color, corner_radius=0.15)
    add_text_box(slide, x + Inches(2.1), Inches(2.22), Inches(0.55), Inches(0.55),
                 str(i + 1), font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow connectors (text-based)
for i in range(3):
    x = Inches(3.2) + i * Inches(3.15)
    add_text_box(slide, x, Inches(3.3), Inches(0.5), Inches(0.5),
                 "→", font_size=28, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom quote
add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.5), Inches(0.8),
             '"Shift-left security for AI agent infrastructure — catch misconfigurations before they become breaches."',
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 4: HOW IT WORKS — Architecture
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "How It Works", "Architecture & Data Flow")

# Architecture boxes
components = [
    (Inches(0.8), Inches(2.3), Inches(2.5), Inches(1.5), "🖥️ Dashboard", "Next.js 14\nAuto-refresh 15s\nPort :3000", PURPLE),
    (Inches(4.0), Inches(2.3), Inches(5.0), Inches(1.5), "⚙️ Governance Controller", "Go API Server :8090  |  Scoring Engine  |  🧠 AI Agent  |  Discovery", CYAN),
    (Inches(10.0), Inches(2.3), Inches(2.5), Inches(1.5), "☸ K8s API", "list / watch\nAll namespaces", RGBColor(0x64, 0x74, 0x8B)),
]

for x, y, w, h, title, desc, color in components:
    add_shape(slide, x, y, w, h, BG_CARD, corner_radius=0.04)
    add_shape(slide, x, y, w, Inches(0.04), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), w - Inches(0.4), Inches(0.8),
                 desc, font_size=11, color=GRAY)

# Arrows
add_text_box(slide, Inches(3.3), Inches(2.7), Inches(0.7), Inches(0.5),
             "←→", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.1), Inches(2.7), Inches(0.9), Inches(0.5),
             "←→", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Discovered Resources row
resources = [
    ("🛡️", "AgentGateway\nBackends, Policies", AMBER),
    ("🤖", "Kagent\nAgents, MCPServers", GREEN),
    ("🌐", "Gateway API\nGateways, HTTPRoutes", RGBColor(0x8B, 0x5C, 0xF6)),
    ("📋", "Governance CRDs\nPolicy, Evaluation", RED),
    ("🧠", "LLM Providers\nGemini, Ollama", PINK),
]

for i, (icon, label, color) in enumerate(resources):
    x = Inches(0.5) + i * Inches(2.55)
    add_shape(slide, x, Inches(4.5), Inches(2.3), Inches(1.2), BG_CARD, corner_radius=0.04)
    add_text_box(slide, x + Inches(0.15), Inches(4.55), Inches(0.5), Inches(0.5),
                 icon, font_size=22, color=color)
    add_text_box(slide, x + Inches(0.6), Inches(4.55), Inches(1.6), Inches(1.1),
                 label, font_size=11, color=LIGHT_GRAY)

# Data flow steps
flow_steps = [
    "1. Controller discovers MCP resources every 30s",
    "2. Reads MCPGovernancePolicy CRD for rules",
    "3. Evaluator scores across 8 categories (0–100)",
    "4. AI Agent sends state to LLM for deep analysis",
    "5. Results exposed via REST API + written to CRD",
    "6. Dashboard polls API every 15s for real-time view",
]
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.4),
             "Data Flow:", font_size=14, color=WHITE, bold=True)
flow_text = "   →   ".join(flow_steps)
add_text_box(slide, Inches(0.8), Inches(6.35), Inches(12), Inches(0.5),
             flow_text, font_size=10, color=GRAY)


# ══════════════════════════════════════════════════════
# SLIDE 5: SCORING MODEL
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Scoring Model", "Policy-driven 0–100 governance score across 8 categories")

# Score gauge visual
gauge = add_shape(slide, Inches(0.8), Inches(2.2), Inches(3.5), Inches(4.5), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(0.8), Inches(2.5), Inches(3.5), Inches(0.5),
             "Score Gauge", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Grade scale
grades = [
    ("A", "90–100", "Compliant", GREEN),
    ("B", "70–89", "Warning", CYAN),
    ("C", "50–69", "NonCompliant", AMBER),
    ("D", "30–49", "NonCompliant", RGBColor(0xF9, 0x73, 0x16)),
    ("F", "0–29", "Critical", RED),
]
for i, (grade, range_str, phase, color) in enumerate(grades):
    y = Inches(3.1) + i * Inches(0.65)
    add_shape(slide, Inches(1.1), y, Inches(0.5), Inches(0.5), color, corner_radius=0.1)
    add_text_box(slide, Inches(1.1), y + Inches(0.02), Inches(0.5), Inches(0.5),
                 grade, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(1), Inches(0.4),
                 range_str, font_size=14, color=LIGHT_GRAY)
    add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(1.3), Inches(0.4),
                 phase, font_size=12, color=color)

# Category weights
add_text_box(slide, Inches(4.8), Inches(2.2), Inches(4), Inches(0.4),
             "Category Weights", font_size=18, color=WHITE, bold=True)

categories = [
    ("AgentGateway Compliance", 25, "Critical"),
    ("Authentication (JWT)", 20, "High"),
    ("Authorization (RBAC)", 15, "High"),
    ("TLS Encryption", 10, "High"),
    ("CORS Policy", 10, "Medium"),
    ("Prompt Guard", 10, "Medium"),
    ("Rate Limiting", 5, "Medium"),
    ("Tool Scope", 5, "Warning"),
]

max_bar = Inches(2.5)
for i, (cat, weight, sev) in enumerate(categories):
    y = Inches(2.75) + i * Inches(0.52)
    add_text_box(slide, Inches(4.8), y, Inches(2.5), Inches(0.35),
                 cat, font_size=11, color=LIGHT_GRAY)
    bar_w = int(max_bar * weight / 25)
    bar_color = GREEN if weight >= 20 else CYAN if weight >= 10 else AMBER
    add_shape(slide, Inches(7.3), y + Inches(0.05), bar_w, Inches(0.25), bar_color, corner_radius=0.15)
    add_text_box(slide, Inches(7.3), y + Inches(0.02), bar_w, Inches(0.3),
                 f"{weight}%", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Formula
add_shape(slide, Inches(9.5), Inches(2.2), Inches(3.5), Inches(2.3), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(9.5), Inches(2.3), Inches(3.5), Inches(0.4),
             "Scoring Formula", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.7), Inches(2.8), Inches(3.1), Inches(0.5),
             "Category Score =", font_size=13, color=CYAN, bold=True)
add_text_box(slide, Inches(9.7), Inches(3.15), Inches(3.1), Inches(0.4),
             "max(0, 100 - Σ penalties)", font_size=12, color=LIGHT_GRAY)
add_text_box(slide, Inches(9.7), Inches(3.6), Inches(3.1), Inches(0.5),
             "Cluster Score =", font_size=13, color=CYAN, bold=True)
add_text_box(slide, Inches(9.7), Inches(3.95), Inches(3.1), Inches(0.4),
             "Σ(score × weight) / Σ weights", font_size=12, color=LIGHT_GRAY)

# Severity penalties
add_shape(slide, Inches(9.5), Inches(4.8), Inches(3.5), Inches(1.9), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(9.5), Inches(4.9), Inches(3.5), Inches(0.4),
             "Severity Penalties", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
penalties = [("Critical", "−40 pts", RED), ("High", "−25 pts", AMBER),
             ("Medium", "−15 pts", RGBColor(0xFB, 0xBF, 0x24)), ("Low", "−5 pts", GRAY)]
for i, (sev, pts, color) in enumerate(penalties):
    y = Inches(5.4) + i * Inches(0.32)
    add_text_box(slide, Inches(9.9), y, Inches(1.5), Inches(0.3), sev, font_size=12, color=color, bold=True)
    add_text_box(slide, Inches(11.5), y, Inches(1.2), Inches(0.3), pts, font_size=12, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════
# SLIDE 6: AI-POWERED SCORING
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "🧠 AI-Powered Governance Scoring",
                   "Deep risk analysis powered by Google Gemini or local Ollama")

# Left: What AI provides
add_shape(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.8), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.4),
             "What the AI Agent Provides", font_size=18, color=PURPLE_LIGHT, bold=True)

ai_features = [
    ("🎯  AI Score", "Independent 0–100 score with grade, generated by the LLM"),
    ("💬  Reasoning", "Human-readable explanation of why the AI assigned its score"),
    ("⚠️  Risk Analysis", "Categorized risks with severity, description, and impact"),
    ("💡  Suggestions", "Prioritized, actionable remediation steps"),
    ("⚖️  Comparison", "Side-by-side AI score vs algorithmic score"),
]

for i, (title, desc) in enumerate(ai_features):
    y = Inches(3.0) + i * Inches(0.75)
    add_text_box(slide, Inches(1.1), y, Inches(2.5), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.1), y + Inches(0.3), Inches(5), Inches(0.35),
                 desc, font_size=11, color=GRAY)

# Right: Provider cards
providers = [
    ("Google Gemini", "gemini-2.5-flash", "Cloud API\n20 req/day (free tier)\nGOOGLE_API_KEY secret", PURPLE_LIGHT),
    ("Ollama (Local)", "llama3.1, qwen2.5, etc.", "Air-gapped / private\nOpenAI-compatible API\nNo API key needed", GREEN),
]

for i, (name, model, desc, color) in enumerate(providers):
    y = Inches(2.2) + i * Inches(2.5)
    add_shape(slide, Inches(6.8), y, Inches(2.8), Inches(2.2), BG_CARD, corner_radius=0.05)
    add_shape(slide, Inches(6.8), y, Inches(0.06), Inches(2.2), color)
    add_text_box(slide, Inches(7.1), y + Inches(0.2), Inches(2.4), Inches(0.35),
                 name, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(7.1), y + Inches(0.6), Inches(2.4), Inches(0.3),
                 f"Model: {model}", font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, Inches(7.1), y + Inches(0.95), Inches(2.4), Inches(1),
                 desc, font_size=10, color=GRAY)

# Rate limiting card
add_shape(slide, Inches(10.0), Inches(2.2), Inches(2.8), Inches(4.8), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(10.0), Inches(2.3), Inches(2.8), Inches(0.4),
             "⏱️ Smart Rate Limiting", font_size=14, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)

rate_items = [
    "Configurable scan\ninterval (default: 5m)",
    "Exponential backoff\n5m → 10m → 20m → 30m",
    "Auto-reset on success",
    "Dashboard refresh\nbypasses rate limit",
    "Pause/Resume toggle\nwithout CR changes",
]
for i, item in enumerate(rate_items):
    y = Inches(2.85) + i * Inches(0.8)
    add_text_box(slide, Inches(10.2), y, Inches(2.5), Inches(0.7),
                 f"▸ {item}", font_size=10, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════
# SLIDE 7: FEATURES — Dashboard & Controls
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Enterprise Dashboard", "Real-time visibility into your MCP security posture")

# Dashboard feature cards
dash_features = [
    ("📊", "Score Gauge", "Animated 0–100 dial with\ngrade (A–F) and phase", PURPLE),
    ("🧠", "AI Score Card", "AI analysis with risks,\nsuggestions, refresh/pause", PURPLE_LIGHT),
    ("📈", "Trend Chart", "Historical score and\nfinding count over time", CYAN),
    ("🔍", "Findings Table", "Filterable by severity,\ncategory, resource", AMBER),
    ("📦", "Resource Cards", "Agents, MCPServers,\nGateways, Endpoints", GREEN),
    ("📋", "Category Breakdown", "Per-category scores\nwith weight visualization", RGBColor(0x8B, 0x5C, 0xF6)),
    ("💡", "Score Explainer", "Human-readable scoring\nexplanation", PINK),
    ("🗂️", "Resource Inventory", "Per-resource drill-down\nwith individual scores", RGBColor(0x64, 0x74, 0x8B)),
]

for i, (icon, title, desc, color) in enumerate(dash_features):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + col * Inches(3.2)
    y = Inches(2.2) + row * Inches(2.4)
    add_icon_card(slide, x, y, Inches(2.9), Inches(2.0), icon, title, desc, color)

# Dashboard tabs
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
             "3 Tabs:   📊 Overview   |   📦 Resources   |   🔍 Findings        •  Auto-refresh every 15 seconds",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# SLIDE 8: FEATURES — CRD & Configuration
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Configuration & Custom Resources",
                   "Everything is driven by Kubernetes CRDs — no external config needed")

# MCPGovernancePolicy card
add_shape(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.8), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5.5), Inches(0.4),
             "📋 MCPGovernancePolicy", font_size=20, color=PURPLE, bold=True)
add_text_box(slide, Inches(1.1), Inches(2.85), Inches(5.5), Inches(0.35),
             "Cluster-scoped  •  Short name: mgp", font_size=12, color=GRAY)

policy_items = [
    "require* flags — toggle each governance category",
    "scoringWeights — customize category importance",
    "severityPenalties — tune penalty per severity",
    "targetNamespaces / excludeNamespaces — scope control",
    "maxToolsWarning / maxToolsCritical — tool thresholds",
    "aiAgent block — provider, model, scan interval, on/off",
    "Status subresource — auto-populated with score & phase",
]
add_bullet_list(slide, Inches(1.1), Inches(3.3), Inches(5.5), Inches(3.5),
                policy_items, font_size=13)

# GovernanceEvaluation card
add_shape(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(4.8), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5.5), Inches(0.4),
             "📊 GovernanceEvaluation", font_size=20, color=CYAN, bold=True)
add_text_box(slide, Inches(7.3), Inches(2.85), Inches(5.5), Inches(0.35),
             "Cluster-scoped  •  Short names: ge, goveval", font_size=12, color=GRAY)

eval_items = [
    "policyRef — references the governance policy",
    "evaluationScope — cluster / namespace / resource",
    "status.score — overall governance score (0–100)",
    "status.findings — full list with remediation steps",
    "status.scoreBreakdown — per-category scores",
    "status.resourceSummary — discovered resource counts",
    "status.namespaceScores — per-namespace breakdown",
]
add_bullet_list(slide, Inches(7.3), Inches(3.3), Inches(5.5), Inches(3.5),
                eval_items, font_size=13, bullet_color=CYAN)


# ══════════════════════════════════════════════════════
# SLIDE 9: API & DEPLOYMENT
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "REST API & Deployment", "12 API endpoints + multiple deployment options")

# API endpoints
add_shape(slide, Inches(0.8), Inches(2.2), Inches(6.5), Inches(4.8), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(1.1), Inches(2.35), Inches(6), Inches(0.4),
             "📜 REST API Endpoints (port 8090)", font_size=16, color=WHITE, bold=True)

endpoints = [
    ("GET", "/api/health", "Health check + version"),
    ("GET", "/api/governance/score", "Score, grade, phase, categories"),
    ("GET", "/api/governance/findings", "All findings + severity breakdown"),
    ("GET", "/api/governance/resources", "Resource inventory summary"),
    ("GET", "/api/governance/resources/detail", "Per-resource scores & findings"),
    ("GET", "/api/governance/breakdown", "Category score breakdown"),
    ("GET", "/api/governance/trends", "Historical score data points"),
    ("GET", "/api/governance/ai-score", "AI score, risks, suggestions"),
    ("POST", "/api/governance/ai-score/refresh", "Trigger immediate AI eval"),
    ("POST", "/api/governance/ai-score/toggle", "Pause/resume AI scanning"),
]

for i, (method, path, desc) in enumerate(endpoints):
    y = Inches(2.85) + i * Inches(0.38)
    method_color = GREEN if method == "GET" else AMBER
    add_text_box(slide, Inches(1.1), y, Inches(0.6), Inches(0.3),
                 method, font_size=10, color=method_color, bold=True)
    add_text_box(slide, Inches(1.7), y, Inches(2.8), Inches(0.3),
                 path, font_size=10, color=CYAN)
    add_text_box(slide, Inches(4.5), y, Inches(2.5), Inches(0.3),
                 desc, font_size=10, color=GRAY)

# Deployment options
add_shape(slide, Inches(7.7), Inches(2.2), Inches(5.0), Inches(4.8), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(8.0), Inches(2.35), Inches(4.5), Inches(0.4),
             "🚀 Deployment Options", font_size=16, color=WHITE, bold=True)

deploy_options = [
    ("Kind (Local Dev)", "make create-cluster && make all\nOne-command local setup", GREEN),
    ("Helm Chart", "helm install mcp-governance\n./charts/mcp-governance", PURPLE),
    ("Raw Manifests", "kubectl apply -f deploy/\nFull control over resources", CYAN),
    ("Production (Registry)", "Push to GHCR / Docker Hub\nimagePullPolicy: Always", AMBER),
]

for i, (title, desc, color) in enumerate(deploy_options):
    y = Inches(2.9) + i * Inches(1.1)
    add_shape(slide, Inches(7.7), y, Inches(0.06), Inches(0.9), color)
    add_text_box(slide, Inches(8.1), y + Inches(0.05), Inches(4.3), Inches(0.3),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(8.1), y + Inches(0.35), Inches(4.3), Inches(0.55),
                 desc, font_size=10, color=GRAY)


# ══════════════════════════════════════════════════════
# SLIDE 10: CI/CD & TESTING
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "CI/CD & Testing", "Automated pipelines with comprehensive test coverage")

# CI Pipeline
add_shape(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(2.3), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(1.1), Inches(2.35), Inches(5.5), Inches(0.4),
             "🔄 CI Pipeline (Push / PR to main)", font_size=16, color=GREEN, bold=True)
ci_items = ["go vet → go test → go build (controller)", "npm ci → npm run build (dashboard)", "helm lint (chart validation)"]
add_bullet_list(slide, Inches(1.1), Inches(2.85), Inches(5.5), Inches(1.5), ci_items, font_size=13, bullet_color=GREEN)

# Release Pipeline
add_shape(slide, Inches(7.0), Inches(2.2), Inches(5.8), Inches(2.3), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(7.3), Inches(2.35), Inches(5.5), Inches(0.4),
             "🏷️ Release Pipeline (Push v* tag)", font_size=16, color=PURPLE_LIGHT, bold=True)
release_items = ["Multi-arch images (amd64 + arm64) → GHCR", "Helm chart → GHCR OCI registry", "GitHub Release with install instructions"]
add_bullet_list(slide, Inches(7.3), Inches(2.85), Inches(5.5), Inches(1.5), release_items, font_size=13, bullet_color=PURPLE_LIGHT)

# Test coverage
add_shape(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(2.3), BG_CARD, corner_radius=0.05)
add_text_box(slide, Inches(1.1), Inches(4.95), Inches(11.5), Inches(0.4),
             "🧪 Test Coverage", font_size=18, color=CYAN, bold=True)

test_packages = [
    ("pkg/evaluator", "85+ tests", "All 9 governance checks, scoring engine, penalties, weights, edge cases"),
    ("pkg/discovery", "17 tests", "Nested object traversal helpers (getNestedMap, String, Int, Slice)"),
    ("cmd/api", "30+ tests", "All HTTP handlers, helpers, CORS, trend recording, JSON responses"),
]

for i, (pkg, count, desc) in enumerate(test_packages):
    x = Inches(1.1) + i * Inches(3.9)
    add_text_box(slide, x, Inches(5.5), Inches(3.5), Inches(0.3),
                 pkg, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x, Inches(5.85), Inches(3.5), Inches(0.3),
                 count, font_size=13, color=CYAN, bold=True)
    add_text_box(slide, x, Inches(6.15), Inches(3.5), Inches(0.7),
                 desc, font_size=10, color=GRAY)


# ══════════════════════════════════════════════════════
# SLIDE 11: CLOSING
# ══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Decorative circles
for cx, cy, r, c in [
    (Inches(10), Inches(0.5), Inches(3.5), PURPLE),
    (Inches(0.5), Inches(4.5), Inches(3), PURPLE_LIGHT),
]:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    circle.fill.solid()
    circle.fill.fore_color.rgb = c
    circle.line.fill.background()

add_text_box(slide, Inches(2), Inches(1.5), Inches(9.5), Inches(1),
             "MCP Governance", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

badge = add_shape(slide, Inches(5.15), Inches(2.6), Inches(3), Inches(0.5), PURPLE_LIGHT, corner_radius=0.15)
add_text_box(slide, Inches(5.15), Inches(2.62), Inches(3), Inches(0.5),
             "🧠  AI-Powered", font_size=22, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.4), Inches(9.5), Inches(0.8),
             "Secure your AI agent infrastructure with\npolicy-driven governance and AI-powered insights.",
             font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Key stats
stats = [
    ("8", "Governance\nCategories"),
    ("0–100", "Weighted\nScore"),
    ("2", "LLM\nProviders"),
    ("12", "REST API\nEndpoints"),
    ("130+", "Unit\nTests"),
]

for i, (num, label) in enumerate(stats):
    x = Inches(1.3) + i * Inches(2.2)
    add_shape(slide, x, Inches(4.6), Inches(1.8), Inches(1.3), BG_CARD, corner_radius=0.06)
    add_text_box(slide, x, Inches(4.7), Inches(1.8), Inches(0.6),
                 num, font_size=32, color=PURPLE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.2), Inches(1.8), Inches(0.6),
                 label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Links
add_text_box(slide, Inches(2), Inches(6.3), Inches(9.5), Inches(0.5),
             "🔗 github.com/techwithhuz/mcp-security-governance    |    📄 MIT License",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── SAVE ────────────────────────────────────────────
output_path = "MCP-Governance-Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
